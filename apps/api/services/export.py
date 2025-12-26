import io
import csv
import zipfile
import json
import os
import tempfile
from typing import List
from datetime import datetime
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use('Agg') # Non-interactive backend
import matplotlib.pyplot as plt

from apps.api.db import models

async def generate_csv_export(workspace_id: str, db: AsyncSession) -> io.BytesIO:
    # Prepare Zip
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
        # 1. Sources
        sources_stmt = select(models.Source).where(models.Source.workspace_id == workspace_id)
        sources = (await db.execute(sources_stmt)).scalars().all()
        
        s_csv = io.StringIO()
        writer = csv.writer(s_csv)
        writer.writerow(["id", "type", "title", "url", "status", "created_at"])
        for s in sources:
            writer.writerow([s.id, s.type, s.title, s.url, s.status, s.created_at])
        zip_file.writestr("sources.csv", s_csv.getvalue())

        # 2. Scorecard Results
        # Get all scorecards for workspace
        sc_stmt = select(models.Scorecard).where(models.Scorecard.workspace_id == workspace_id)
        scorecards = (await db.execute(sc_stmt)).scalars().all()
        sc_ids = [sc.id for sc in scorecards]
        
        if sc_ids:
            scr_stmt = select(models.ScorecardResult).where(models.ScorecardResult.scorecard_id.in_(sc_ids))
            results = (await db.execute(scr_stmt)).scalars().all()
            
            r_csv = io.StringIO()
            writer = csv.writer(r_csv)
            writer.writerow(["scorecard_id", "brand", "overall_score", "factor_breakdown"])
            for r in results:
                writer.writerow([r.scorecard_id, r.brand, r.results.get('overall'), json.dumps(r.results.get('factors'))])
            zip_file.writestr("scorecard_results.csv", r_csv.getvalue())

        # 3. Insights (Claims, Sentiment, etc.)
        i_stmt = select(models.Insight).where(models.Insight.workspace_id == workspace_id)
        insights = (await db.execute(i_stmt)).scalars().all()
        
        # Flatten insights
        # Claims
        claims_insight = next((i for i in insights if i.kind == 'claims'), None)
        if claims_insight and claims_insight.metrics:
            c_csv = io.StringIO()
            writer = csv.writer(c_csv)
            writer.writerow(["claim", "count"])
            for claim, count in claims_insight.metrics.items():
                writer.writerow([claim, count])
            zip_file.writestr("claims.csv", c_csv.getvalue())
            
        # Stats
        stats_insight = next((i for i in insights if i.kind == 'stats'), None)
        if stats_insight and stats_insight.metrics:
             # Just dump json for stats
             zip_file.writestr("stats.json", json.dumps(stats_insight.metrics, indent=2))

    zip_buffer.seek(0)
    return zip_buffer

async def generate_pptx_export(workspace_id: str, db: AsyncSession) -> str:
    # Fetch Data
    # Workspace info
    ws = await db.get(models.Workspace, workspace_id)
    ws_name = ws.name if ws else "Workspace"
    
    # Insights
    i_stmt = select(models.Insight).where(models.Insight.workspace_id == workspace_id)
    insights = (await db.execute(i_stmt)).scalars().all()
    
    stats = next((i.metrics for i in insights if i.kind == 'stats'), {}) or {}
    claims = next((i.metrics for i in insights if i.kind == 'claims'), {}) or {}
    trends = next((i.metrics for i in insights if i.kind == 'trends'), {}) or {}
    themes = [i for i in insights if i.kind == 'theme']

    # Initialize Presentation
    prs = Presentation()

    # Helper to add title slide
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    # Helper to add content slide
    def add_content_slide(title):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        return slide

    # Slide 1: Title
    add_title_slide(f"InsightHub Report: {ws_name}", f"Generated on {datetime.now().strftime('%Y-%m-%d')}")

    # Slide 2: Executive Summary
    slide = add_content_slide("Executive Summary")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = f"Overview of analysis for {ws_name}."
    p = tf.add_paragraph()
    p.text = f"Total Documents Processed: {stats.get('total_documents', 'N/A')}"
    p = tf.add_paragraph()
    p.text = f"Overall Sentiment: {stats.get('average_sentiment', 'N/A')}"
    p = tf.add_paragraph()
    p.text = f"Key Themes Identified: {len(themes)}"

    # Slide 3: Claims Landscape (Chart)
    slide = add_content_slide("Claims Landscape")
    if claims:
        # Generate Chart
        fig, ax = plt.subplots(figsize=(6, 4))
        # Top 10 claims
        top_claims = sorted(claims.items(), key=lambda x: x[1], reverse=True)[:10]
        labels = [x[0] for x in top_claims]
        values = [x[1] for x in top_claims]
        
        ax.barh(labels, values, color='skyblue')
        ax.set_xlabel('Count')
        ax.set_title('Top Claims Mentions')
        plt.tight_layout()
        
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            plt.savefig(tmp.name)
            img_path = tmp.name
        plt.close()

        # Insert Image
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))
        os.unlink(img_path)
    else:
        slide.shapes.placeholders[1].text = "No claims data available."

    # Slide 4: Sentiment Overview
    slide = add_content_slide("Sentiment Drivers")
    dist = stats.get('sentiment_distribution', {})
    if dist:
        fig, ax = plt.subplots(figsize=(5, 5))
        ax.pie(dist.values(), labels=dist.keys(), autopct='%1.1f%%', colors=['#ff9999','#66b3ff','#99ff99'])
        ax.set_title('Sentiment Distribution')
        plt.tight_layout()
        
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            plt.savefig(tmp.name)
            img_path = tmp.name
        plt.close()
        
        slide.shapes.add_picture(img_path, Inches(2.5), Inches(2), height=Inches(5))
        os.unlink(img_path)
    else:
        slide.shapes.placeholders[1].text = "No sentiment data available."

    # Slide 5: Themes
    slide = add_content_slide("Emerging Themes")
    tf = slide.shapes.placeholders[1].text_frame
    for theme in themes[:5]: # Top 5 themes
        title = theme.title
        count = theme.metrics.get('count', 0)
        p = tf.add_paragraph()
        p.text = f"{title} (Count: {count})"
        p.level = 0
        p = tf.add_paragraph()
        p.text = theme.summary
        p.level = 1

    # Slide 6: Trends (if any)
    slide = add_content_slide("Trends Analysis")
    if trends:
        # Plot lines for each brand
        fig, ax = plt.subplots(figsize=(8, 4))
        for brand, data in trends.items():
            # sort by date key
            dates = sorted(data.keys())
            values = [data[d] for d in dates]
            ax.plot(dates, values, marker='o', label=brand)
        
        ax.legend()
        ax.set_title("Average Rating Over Time")
        plt.xticks(rotation=45)
        plt.tight_layout()

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            plt.savefig(tmp.name)
            img_path = tmp.name
        plt.close()
        
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))
        os.unlink(img_path)
    else:
        slide.shapes.placeholders[1].text = "No trend data available."

    # Scorecard Results
    scr_stmt = select(models.ScorecardResult).join(models.Scorecard).where(models.Scorecard.workspace_id == workspace_id)
    scorecard_results = (await db.execute(scr_stmt)).scalars().all()

    # Slide 7: Scorecard Comparison
    slide = add_content_slide("Scorecard Comparison")
    if scorecard_results:
        # Bar chart of overall scores by brand
        brands = [r.brand for r in scorecard_results]
        scores = [r.results.get('overall', 0) for r in scorecard_results]
        
        fig, ax = plt.subplots(figsize=(6, 4))
        ax.bar(brands, scores, color='lightgreen')
        ax.set_ylim(0, 100)
        ax.set_ylabel('Score')
        ax.set_title('Overall Brand Scores')
        plt.tight_layout()

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            plt.savefig(tmp.name)
            img_path = tmp.name
        plt.close()
        
        slide.shapes.add_picture(img_path, Inches(2), Inches(2), width=Inches(6))
        os.unlink(img_path)
    else:
         slide.shapes.placeholders[1].text = "No scorecard results available."


    # Slide 8: Recommendations (Template)
    slide = add_content_slide("Recommendations")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Based on the analysis, we recommend:"
    p = tf.add_paragraph()
    p.text = "[Placeholder for AI-generated recommendations based on whitelist gaps]"
    p = tf.add_paragraph()
    p.text = "[Placeholder for product improvement suggestions]"

    # Slide 9: Appendix (Sources)
    slide = add_content_slide("Appendix: Data Sources")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Key sources analyzed:"
    
    # List first 10 sources
    sources_stmt = select(models.Source).where(models.Source.workspace_id == workspace_id).limit(10)
    sources = (await db.execute(sources_stmt)).scalars().all()
    for s in sources:
        p = tf.add_paragraph()
        p.text = f"- {s.title} ({s.type})"
        p.level = 0


    # Slide 10: End
    add_title_slide("Thank You", "InsightHub Generated Report")

    # Save
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return tmp.name
